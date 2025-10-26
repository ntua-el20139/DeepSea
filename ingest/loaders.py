from typing import Iterator, Tuple, Optional, List, Dict, Any, TYPE_CHECKING
from pathlib import Path
import io
from PIL import Image
import pdfplumber
from pypdf import PdfReader
from pptx import Presentation
from docx import Document as Docx

from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.accelerator_options import AcceleratorOptions
from docling.datamodel.base_models import ConversionStatus, InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions

from ingest.docling_ocr import ocr_image

print("[loaders] module loaded")

if TYPE_CHECKING:
    from docling_core.types.doc.document import DoclingDocument

_DOC_TABLE_CONVERTER: Optional[DocumentConverter] = None


def _pptx_table_to_markdown(shape) -> Optional[str]:
    """
    Return a markdown representation of the table contained in `shape`
    using the same restrictions as `_rows_to_markdown`.
    """
    try:
        table = shape.table
    except (AttributeError, ValueError):
        return None

    # Extract and clean rows
    rows: List[List[str]] = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            # Normalize to string, strip, squash newlines (PPTX often has them)
            txt = (cell.text or "").replace("\r", "\n").replace("\n", " ").strip()
            cells.append(txt)
        rows.append(cells)

    # Delegate validation/formatting
    return _rows_to_markdown(rows)


def _rows_to_markdown(rows: List[List[str]]) -> Optional[str]:
    """
    Convert a 2D list of strings into a simple GitHub-style markdown table.
    Pads ragged rows to the max column width for consistent formatting.
    """
    if not rows or len(rows) < 2:
        return None
    # Clean input: strip whitespace and normalize None → ""
    cleaned_rows: List[List[str]] = [
        [("" if cell is None else str(cell)).strip() for cell in row]
        for row in rows
    ]
    header = cleaned_rows[0]
    data_rows = cleaned_rows[1:]
    # Header must not be empty or contain empty cells
    if not header or any(h == "" for h in header):
        return None
    header_len = len(header)

    valid_rows = []
    for row in data_rows:
        # Reject if too many columns
        if len(row) > header_len:
            return None
        # Pad if too few
        if len(row) < header_len:
            row = row + [""] * (header_len - len(row))
        valid_rows.append(row)

    # Reject if there's no valid data row
    if not valid_rows:
        return None

    # Build markdown table
    header_line = "| " + " | ".join(header) + " |"
    separator = "| " + " | ".join(["---"] * header_len) + " |"
    body_lines = ["| " + " | ".join(r) + " |" for r in valid_rows]

    return "\n".join([header_line, separator, *body_lines])


def _pptx_image(shape) -> Optional[Image.Image]:
    """
    Return a PIL ``Image`` if ``shape`` embeds a picture, otherwise ``None``.
    """
    try:
        blob = shape.image.blob
    except (AttributeError, ValueError):
        return None

    try:
        return Image.open(io.BytesIO(blob)).convert("RGB")
    except Exception:
        return None


def _docx_table_to_md(tbl) -> Optional[str]:
    """
    Convert a python-docx table to GitHub-style markdown with the same
    restrictions as `_rows_to_markdown`.
    """
    try:
        rows: List[List[str]] = []
        for row in tbl.rows:
            cells = []
            for cell in row.cells:
                txt = (cell.text or "").replace("\r", "\n").replace("\n", " ").strip()
                cells.append(txt)
            rows.append(cells)
    except Exception:
        return None

    return _rows_to_markdown(rows)


def load_pdf_text(path: str) -> Iterator[Tuple[int, str]]:
    print(f"[loaders] load_pdf_text: {path}")
    reader = PdfReader(path)
    for i, page in enumerate(reader.pages):
        yield (i + 1), (page.extract_text() or "")


def _get_docling_converter() -> DocumentConverter:
    global _DOC_TABLE_CONVERTER
    if _DOC_TABLE_CONVERTER is None:
        pdf_option = PdfFormatOption(
            pipeline_options=PdfPipelineOptions(
                do_ocr=False,
                accelerator_options=AcceleratorOptions(device="cpu"),
            )
        )
        _DOC_TABLE_CONVERTER = DocumentConverter(
            allowed_formats=[InputFormat.PDF],
            format_options={InputFormat.PDF: pdf_option},
        )
    return _DOC_TABLE_CONVERTER


def _clean_docling_cell_text(cell: Any, doc: "DoclingDocument") -> str:
    try:
        text = cell._get_text(doc=doc)
    except Exception:
        text = getattr(cell, "text", "")
    text = "" if text is None else str(text)
    text = text.replace("\r", " ").replace("\n", " ")
    return " ".join(text.split())


def _docling_table_to_rows(table: Any, doc: "DoclingDocument") -> Optional[List[List[str]]]:
    data = getattr(table, "data", None)
    grid = getattr(data, "grid", None) if data is not None else None
    if not grid:
        return None

    header_rows: List[List[str]] = []
    data_rows: List[List[str]] = []

    for row in grid:
        cleaned_row = [_clean_docling_cell_text(cell, doc) for cell in row]
        if not any(cleaned_row):
            continue
        row_is_header = any(getattr(cell, "column_header", False) for cell in row)
        if row_is_header and not data_rows:
            header_rows.append(cleaned_row)
        else:
            data_rows.append(cleaned_row)

    if not header_rows and not data_rows:
        return None

    if header_rows:
        num_cols = len(header_rows[0])
        header = []
        for col_idx in range(num_cols):
            parts = [row[col_idx] for row in header_rows if row[col_idx]]
            header.append(" ".join(parts).strip())
    else:
        header = data_rows[0]
        data_rows = data_rows[1:]

    if not data_rows:
        return None

    header = [col if col else f"Column {idx + 1}" for idx, col in enumerate(header)]

    return [header, *data_rows]


def load_pdf_tables(path: str) -> Dict[int, List[str]]:
    """
    Extract tables from a PDF, returning markdown strings keyed by 1-based page number.
    """
    print(f"[loaders] load_pdf_tables: {path}")
    tables_by_page: Dict[int, List[str]] = {}
    try:
        converter = _get_docling_converter()
        conv_res = converter.convert(path)
    except Exception as e:
        print(f"[loaders] docling table extraction failed ({e}); skipping")
        return tables_by_page

    if conv_res.status not in {ConversionStatus.SUCCESS, ConversionStatus.PARTIAL_SUCCESS}:
        print(
            f"[loaders] docling conversion status {conv_res.status}; skipping tables"
        )
        return tables_by_page

    doc = conv_res.document
    for table in getattr(doc, "tables", []):
        try:
            rows = _docling_table_to_rows(table, doc)
        except Exception as table_err:
            print(f"[loaders] table parsing failed ({table_err}); continuing")
            continue
        if not rows:
            continue
        md = _rows_to_markdown(rows)
        if not md:
            continue
        page_numbers = sorted(
            {
                prov.page_no
                for prov in getattr(table, "prov", [])
                if getattr(prov, "page_no", None) is not None
            }
        )
        if not page_numbers:
            continue
        for page_no in page_numbers:
            tables_by_page.setdefault(page_no, []).append(md)
    return tables_by_page

def _ocr_text_and_confidence(image: Image.Image) -> Tuple[str, Optional[float]]:
    """
    Return OCR text and an average confidence score for the provided image.
    Confidence is expressed on a 0–100 scale to match historical behaviour.
    """
    return ocr_image(image)


def load_pdf_images_ocr(path: str) -> Iterator[Tuple[int, str, Optional[float]]]:
    print(f"[loaders] load_pdf_images_ocr: {path}")
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                im = page.to_image(resolution=220).original
                txt, conf = _ocr_text_and_confidence(im)
                if txt.strip() and conf is not None and conf > 95.0:
                    yield (i + 1), txt, conf
    except Exception as e:
        print(f"[loaders] pdf OCR failed ({e}); skipping")
        return

def load_pptx(path: str) -> Iterator[Tuple[int, str, List[Image.Image], List[str]]]:
    """
    Yield (slide_number, concatenated_text, images[], tables_md[]) for each slide.
    """
    print(f"[loaders] load_pptx: {path}")
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        texts, images, tables_md = [], [], []
        for shape in slide.shapes:
            if hasattr(shape, "text") and getattr(shape, "has_text_frame", False):
                texts.append(shape.text) # type: ignore
            table_md = _pptx_table_to_markdown(shape)
            if table_md:
                tables_md.append(table_md)
            image = _pptx_image(shape)
            if image:
                images.append(image)
        yield (i + 1), "\n".join(texts), images, tables_md

def ocr_images(images: List[Image.Image]) -> Tuple[str, Optional[float]]:
    texts: List[str] = []
    confs: List[float] = []
    for im in images:
        txt, conf = _ocr_text_and_confidence(im)
        if txt.strip() and conf is not None and conf > 95.0:
            texts.append(txt)
            confs.append(conf)
    joined = "\n".join(texts).strip()
    avg_conf = (sum(confs) / len(confs)) if confs else None
    print(f"[loaders] ocr_images: {len(images)} images => {len(joined.split())} words (conf={avg_conf})")
    return joined, avg_conf

def load_plain_text(path: str) -> str:
    print(f"[loaders] load_plain_text: {path}")
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def load_docx_blocks(path: str) -> Iterator[Dict[str, Any]]:
    print(f"[loaders_docx] parsing: {path}")
    d = Docx(path)
    previous = {"status": None, "size": 12, "previous_empty_lines": 0}
    for p in d.paragraphs:
        size = int(p.runs[0].font.size.pt) if p.runs and p.runs[0].font.size else 12
        if p.style and p.style.name and p.style.name.lower().startswith("heading"):     # heading
            try:
                lvl = int(p.style.name.lower().replace("heading", "") or 1)
            except Exception:
                lvl = 1
            if p.text.strip():
                yield {"type": "heading", "level": lvl, "text": p.text}
                previous["status"] = "heading"
                previous["previous_empty_lines"] = 0
                previous["size"] = size

        elif p.text.strip():
            if size > previous["size"] and len(p.runs) == 1:       # bigger font
                if previous["status"] != "bigger_font" and p.text.strip():
                    yield {"type": "bigger_font", "text": p.text}
                previous["status"] = "bigger_font"
                
            elif p.runs[0].font.underline and len(p.runs) == 1:      # underline
                if previous["status"] != "underline" and p.text.strip():
                    yield {"type": "underline", "text": p.text}
                previous["status"] = "underline"
                
            elif p.runs[0].font.bold and len(p.runs) == 1:           # bold
                if previous["status"] != "bold" and p.text.strip():
                    yield {"type": "bold", "text": p.text}
                previous["status"] = "bold"
                
            elif p.runs[0].font.italic and len(p.runs) == 1:         # italic
                if previous["status"] != "italic" and p.text.strip():
                    yield {"type": "italic", "text": p.text}
                previous["status"] = "italic"
                
            elif previous["status"] == "empty" and previous["previous_empty_lines"] >= 2:     # paragraph break
                yield {"type": "paragraph_break", "text": p.text}
                previous["status"] = "paragraph_break"
                
            else:   # normal paragraph
                yield {"type": "paragraph", "text": p.text}
                previous["status"] = "paragraph"
            previous["previous_empty_lines"] = 0
            previous["size"] = size
            
        else:   # empty line
            previous["previous_empty_lines"] += 1
            previous["status"] = "empty"

    for t in d.tables:
        md = _docx_table_to_md(t)
        if md:
            yield {"type": "table", "as_markdown": md}

    for shp in d.inline_shapes:
        try:
            rel_id = shp._inline.graphic.graphicData.pic.blipFill.blip.embed
            img_part = d.part.related_parts[rel_id]
            im = Image.open(io.BytesIO(img_part.blob)).convert("RGB")
            yield {"type": "image", "image": im, "width": im.width, "height": im.height}
        except Exception:
            continue
